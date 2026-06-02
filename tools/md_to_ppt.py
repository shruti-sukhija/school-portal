from pptx import Presentation
from pptx.util import Inches, Pt
import os

MD_PATH = os.path.join(os.path.dirname(__file__), '..', 'docs', 'presentation.md')
OUT_PATH = os.path.join(os.path.dirname(__file__), '..', 'docs', 'presentation.pptx')
IMAGES_DIR = os.path.join(os.path.dirname(__file__), '..', 'School_Images')

# Read markdown and split slides by lines containing only '---'
with open(MD_PATH, 'r', encoding='utf-8') as f:
    md = f.read()

# Normalize separators
parts = []
current = []
for line in md.splitlines():
    if line.strip() == '---':
        parts.append('\n'.join(current).strip())
        current = []
    else:
        current.append(line)
if current:
    parts.append('\n'.join(current).strip())

prs = Presentation()
# set slide width/height defaults if you want (skip)

def add_title_slide(text):
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    # add background image if available
    img_path = os.path.join(IMAGES_DIR, 'School_Building.jpeg')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
    # overlay title box
    left = Inches(0.5)
    top = Inches(1.0)
    width = prs.slide_width - Inches(1)
    height = Inches(2.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = None
    return slide


def add_content_slide(title, bullets, notes=None, image=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
    if image:
        try:
            slide.shapes.add_picture(image, prs.slide_width - Inches(3.0), Inches(1.2), width=Inches(3.0))
        except Exception:
            pass
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# Parse each part
for i, part in enumerate(parts):
    lines = [l.rstrip() for l in part.splitlines() if l.strip()!='']
    if not lines:
        continue
    # title is first line without leading % or #
    title = lines[0].lstrip('% ').lstrip('# ').strip()
    rest = lines[1:]
    # collect bullets: lines starting with '-' or numeric or plain
    bullets = []
    notes = None
    for r in rest:
        if r.startswith('- '):
            bullets.append(r[2:].strip())
        elif r.startswith('* '):
            bullets.append(r[2:].strip())
        elif r.startswith('```'):
            continue
        else:
            # treat as paragraph, may split into sentences
            bullets.extend([s.strip() for s in r.split('. ') if s.strip()])
    # for first slide, use title layout image
    if i == 0:
        s = add_title_slide(title)
        # add subtitle as notes if available
        if bullets:
            s.notes_slide.notes_text_frame.text = '\n'.join(bullets)
    else:
        # try to pick an image from School_Images for this slide
        candidate = None
        img_choices = [
            'School_Building.jpeg', 'Reception.jpeg', 'Robotics/IMG-20260513-WA0054.jpg.jpeg',
            'Smart_Class_Rooms/IMG-20260511-WA0057.jpg.jpeg', 'Sports_and_Culture/IMG-20231205-WA0047.jpg.jpeg'
        ]
        for img in img_choices:
            p = os.path.join(IMAGES_DIR, img)
            if os.path.exists(p):
                candidate = p
                break
        add_content_slide(title, bullets, notes='Speaker notes for ' + title, image=candidate)

prs.save(OUT_PATH)
print('Saved', OUT_PATH)
